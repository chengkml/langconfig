# Copyright (c) 2025 Cade Russell (Ghost Peony)
#
# This source code is licensed under the MIT license found in the
# LICENSE file in the root directory of this source tree.

"""
Presentation Generation Service

Generates presentations from selected workflow artifacts and files.
Supports Google Slides, PDF, and Reveal.js formats.
"""
import os
import io
import base64
import json
import logging
import asyncio
from typing import Optional, List, Dict, Any, Tuple
from datetime import datetime, timezone
from pathlib import Path
import zipfile

from sqlalchemy.orm import Session
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from models.presentation_job import (
    PresentationJob,
    PresentationJobStatus,
    PresentationFormat,
    PresentationTheme
)
from services.oauth_service import google_oauth_service

try:
    import httplib2
    from google_auth_httplib2 import AuthorizedHttp
except ImportError:
    httplib2 = None
    AuthorizedHttp = None

logger = logging.getLogger(__name__)


async def _execute(request):
    """Run a synchronous Google API request off the event loop."""
    return await asyncio.to_thread(request.execute)


class SlideContent:
    """Represents content for a single slide."""
    def __init__(
        self,
        slide_type: str = "content",  # "title", "content", "image", "section"
        title: str = "",
        subtitle: str = "",
        bullets: List[str] = None,
        body_text: str = "",
        image_data: str = None,  # Base64 encoded
        image_mime_type: str = None,
        speaker_notes: str = ""
    ):
        self.slide_type = slide_type
        self.title = title
        self.subtitle = subtitle
        self.bullets = bullets or []
        self.body_text = body_text
        self.image_data = image_data
        self.image_mime_type = image_mime_type
        self.speaker_notes = speaker_notes


class PresentationService:
    """
    Service for generating presentations from workflow artifacts.
    """

    # Maximum time for a single job to complete (4 minutes)
    JOB_TIMEOUT_SECONDS = 240

    # Theme configurations
    THEMES = {
        "default": {
            "background_color": "FFFFFF",
            "title_color": "1A1A1A",
            "text_color": "333333",
            "accent_color": "4F46E5"
        },
        "dark": {
            "background_color": "1A1A2E",
            "title_color": "FFFFFF",
            "text_color": "E0E0E0",
            "accent_color": "818CF8"
        },
        "minimal": {
            "background_color": "FAFAFA",
            "title_color": "111111",
            "text_color": "444444",
            "accent_color": "6366F1"
        }
    }

    def __init__(self):
        self.workspace_base = os.getenv("WORKSPACE_PATH", "./workspace")

    async def create_job(
        self,
        db: Session,
        title: str,
        output_format: str,
        selected_items: List[Dict[str, Any]],
        workflow_id: Optional[int] = None,
        task_id: Optional[int] = None,
        theme: str = "default"
    ) -> PresentationJob:
        """
        Create a new presentation generation job.

        Args:
            db: Database session
            title: Presentation title
            output_format: Output format (google_slides, pdf, revealjs)
            selected_items: List of selected artifacts/files
            workflow_id: Optional workflow ID for context
            task_id: Optional task ID for context
            theme: Visual theme

        Returns:
            The created PresentationJob
        """
        job = PresentationJob(
            title=title,
            output_format=output_format,
            input_items=selected_items,
            workflow_id=workflow_id,
            task_id=task_id,
            theme=theme,
            status=PresentationJobStatus.PENDING.value
        )
        db.add(job)
        db.commit()
        db.refresh(job)

        logger.info(f"Created presentation job {job.id} for format {output_format}")
        return job

    async def process_job(self, db: Session, job_id: int) -> PresentationJob:
        """
        Process a presentation generation job with a timeout.

        Args:
            db: Database session
            job_id: Job ID to process

        Returns:
            The updated job with results
        """
        job = db.query(PresentationJob).filter(PresentationJob.id == job_id).first()
        if not job:
            raise ValueError(f"Job {job_id} not found")

        job.mark_processing()
        db.commit()

        try:
            await asyncio.wait_for(
                self._execute_job(db, job),
                timeout=self.JOB_TIMEOUT_SECONDS
            )
            db.commit()
            logger.info(f"Completed presentation job {job.id}")

        except asyncio.TimeoutError:
            logger.error(f"Presentation job {job.id} timed out after {self.JOB_TIMEOUT_SECONDS}s")
            job.mark_failed(f"Job timed out after {self.JOB_TIMEOUT_SECONDS} seconds")
            db.commit()

        except Exception as e:
            logger.error(f"Failed to process presentation job {job.id}: {e}")
            job.mark_failed(str(e))
            db.commit()

        return job

    async def _execute_job(self, db: Session, job: PresentationJob) -> None:
        """
        Execute the actual presentation generation work.

        Args:
            db: Database session
            job: The job to execute
        """
        # Prepare slide content from input items
        slides = await self._prepare_slides(job.input_items, job.title)

        # Generate based on format
        if job.output_format == PresentationFormat.GOOGLE_SLIDES.value:
            result_url, image_failures = await self._generate_google_slides(db, job.title, slides, job.theme)
            job.mark_completed(result_url=result_url)
            if image_failures:
                job.error_message = f"Completed with {len(image_failures)} image error(s): {'; '.join(image_failures)}"

        elif job.output_format == PresentationFormat.PDF.value:
            file_path = await self._generate_pdf(job.id, job.title, slides, job.theme)
            job.mark_completed(result_file_path=file_path)

        elif job.output_format == PresentationFormat.REVEALJS.value:
            file_path = await self._generate_revealjs(job.id, job.title, slides, job.theme)
            job.mark_completed(result_file_path=file_path)

        else:
            raise ValueError(f"Unknown output format: {job.output_format}")

    async def _prepare_slides(
        self,
        input_items: List[Dict[str, Any]],
        title: str
    ) -> List[SlideContent]:
        """
        Convert input items into structured slide content.

        Args:
            input_items: Selected artifacts and files
            title: Presentation title

        Returns:
            List of SlideContent objects
        """
        slides = []

        # Title slide
        slides.append(SlideContent(
            slide_type="title",
            title=title,
            subtitle=f"Generated on {datetime.now().strftime('%B %d, %Y')}"
        ))

        # Group items by type
        images = []
        text_items = []
        files = []

        for item in input_items:
            item_type = item.get("type", "")
            if item_type == "artifact":
                block = item.get("block", {})
                block_type = block.get("type", "")

                if block_type == "image":
                    images.append({
                        "data": block.get("data", ""),
                        "mimeType": block.get("mimeType", "image/png"),
                        "alt_text": block.get("alt_text", "Generated image")
                    })
                elif block_type == "text":
                    text_items.append(block.get("text", ""))
                elif block_type == "file":
                    files.append({
                        "name": block.get("name", "File"),
                        "text": block.get("text", "")
                    })

            elif item_type == "file":
                files.append({
                    "name": item.get("displayName", item.get("filename", "File")),
                    "path": item.get("filePath", "")
                })

        # Create slides for images (one per slide for prominence)
        for i, img in enumerate(images):
            slides.append(SlideContent(
                slide_type="image",
                title=img.get("alt_text", f"Image {i + 1}"),
                image_data=img.get("data"),
                image_mime_type=img.get("mimeType", "image/png")
            ))

        # Create slides for text content
        if text_items:
            # Split into multiple slides if needed
            combined_text = "\n\n".join(text_items)
            paragraphs = combined_text.split("\n\n")

            # Group paragraphs into slides (roughly 3-4 per slide)
            chunk_size = 4
            for i in range(0, len(paragraphs), chunk_size):
                chunk = paragraphs[i:i + chunk_size]
                slides.append(SlideContent(
                    slide_type="content",
                    title="Key Points" if i == 0 else f"Key Points (continued)",
                    bullets=[p.strip() for p in chunk if p.strip()]
                ))

        # Create slides for files (summary)
        if files:
            file_names = [f.get("name", "Unknown") for f in files]
            slides.append(SlideContent(
                slide_type="content",
                title="Referenced Files",
                bullets=file_names
            ))

        # Closing slide
        slides.append(SlideContent(
            slide_type="section",
            title="Thank You",
            subtitle="Generated with LangConfig"
        ))

        return slides

    async def _generate_pdf(
        self,
        job_id: int,
        title: str,
        slides: List[SlideContent],
        theme: str
    ) -> str:
        """
        Generate a PowerPoint/PDF presentation using python-pptx.

        Args:
            job_id: Job ID for file naming
            title: Presentation title
            slides: List of slide content
            theme: Visual theme

        Returns:
            Path to the generated file
        """
        theme_config = self.THEMES.get(theme, self.THEMES["default"])

        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for slide_content in slides:
            if slide_content.slide_type == "title":
                slide_layout = prs.slide_layouts[6]  # Blank
                slide = prs.slides.add_slide(slide_layout)

                # Title
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
                )
                title_frame = title_box.text_frame
                title_frame.paragraphs[0].text = slide_content.title
                title_frame.paragraphs[0].font.size = Pt(44)
                title_frame.paragraphs[0].font.bold = True
                title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

                # Subtitle
                if slide_content.subtitle:
                    sub_box = slide.shapes.add_textbox(
                        Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8)
                    )
                    sub_frame = sub_box.text_frame
                    sub_frame.paragraphs[0].text = slide_content.subtitle
                    sub_frame.paragraphs[0].font.size = Pt(24)
                    sub_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            elif slide_content.slide_type == "section":
                slide_layout = prs.slide_layouts[6]  # Blank
                slide = prs.slides.add_slide(slide_layout)

                # Centered title
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(3), Inches(12.333), Inches(1.5)
                )
                title_frame = title_box.text_frame
                title_frame.paragraphs[0].text = slide_content.title
                title_frame.paragraphs[0].font.size = Pt(36)
                title_frame.paragraphs[0].font.bold = True
                title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

                if slide_content.subtitle:
                    sub_box = slide.shapes.add_textbox(
                        Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.8)
                    )
                    sub_frame = sub_box.text_frame
                    sub_frame.paragraphs[0].text = slide_content.subtitle
                    sub_frame.paragraphs[0].font.size = Pt(18)
                    sub_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            elif slide_content.slide_type == "content":
                slide_layout = prs.slide_layouts[6]  # Blank
                slide = prs.slides.add_slide(slide_layout)

                # Title
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.5), Inches(12.333), Inches(1)
                )
                title_frame = title_box.text_frame
                title_frame.paragraphs[0].text = slide_content.title
                title_frame.paragraphs[0].font.size = Pt(32)
                title_frame.paragraphs[0].font.bold = True

                # Bullets
                if slide_content.bullets:
                    content_box = slide.shapes.add_textbox(
                        Inches(0.75), Inches(1.7), Inches(11.833), Inches(5)
                    )
                    content_frame = content_box.text_frame
                    content_frame.word_wrap = True

                    for i, bullet in enumerate(slide_content.bullets):
                        if i == 0:
                            p = content_frame.paragraphs[0]
                        else:
                            p = content_frame.add_paragraph()
                        p.text = f"• {bullet}"
                        p.font.size = Pt(18)
                        p.space_before = Pt(12)

            elif slide_content.slide_type == "image":
                slide_layout = prs.slide_layouts[6]  # Blank
                slide = prs.slides.add_slide(slide_layout)

                # Title
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
                )
                title_frame = title_box.text_frame
                title_frame.paragraphs[0].text = slide_content.title
                title_frame.paragraphs[0].font.size = Pt(24)
                title_frame.paragraphs[0].font.bold = True

                # Image
                if slide_content.image_data:
                    try:
                        image_bytes = base64.b64decode(slide_content.image_data)
                        image_stream = io.BytesIO(image_bytes)

                        # Add image centered
                        slide.shapes.add_picture(
                            image_stream,
                            Inches(1.5), Inches(1.3),
                            width=Inches(10.333),
                            height=Inches(5.5)
                        )
                    except Exception as e:
                        logger.warning(f"Failed to add image to slide: {e}")

        # Save to file
        output_dir = Path(self.workspace_base) / "presentations"
        output_dir.mkdir(parents=True, exist_ok=True)

        filename = f"presentation_{job_id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        file_path = output_dir / filename

        prs.save(str(file_path))
        logger.info(f"Generated PPTX at {file_path}")

        return str(file_path)

    async def _generate_revealjs(
        self,
        job_id: int,
        title: str,
        slides: List[SlideContent],
        theme: str
    ) -> str:
        """
        Generate a Reveal.js HTML presentation.

        Args:
            job_id: Job ID for file naming
            title: Presentation title
            slides: List of slide content
            theme: Visual theme

        Returns:
            Path to the generated ZIP file
        """
        theme_config = self.THEMES.get(theme, self.THEMES["default"])

        # Map our themes to Reveal.js themes
        revealjs_theme = {
            "default": "white",
            "dark": "black",
            "minimal": "simple"
        }.get(theme, "white")

        # Generate slide HTML
        slides_html = []
        for slide_content in slides:
            if slide_content.slide_type == "title":
                html = f"""
                <section>
                    <h1>{self._escape_html(slide_content.title)}</h1>
                    <p>{self._escape_html(slide_content.subtitle)}</p>
                </section>
                """
            elif slide_content.slide_type == "section":
                html = f"""
                <section>
                    <h2>{self._escape_html(slide_content.title)}</h2>
                    <p>{self._escape_html(slide_content.subtitle)}</p>
                </section>
                """
            elif slide_content.slide_type == "content":
                bullets_html = "\n".join(
                    f"<li>{self._escape_html(b)}</li>" for b in slide_content.bullets
                )
                html = f"""
                <section>
                    <h2>{self._escape_html(slide_content.title)}</h2>
                    <ul>
                        {bullets_html}
                    </ul>
                </section>
                """
            elif slide_content.slide_type == "image":
                if slide_content.image_data:
                    img_src = f"data:{slide_content.image_mime_type};base64,{slide_content.image_data}"
                    html = f"""
                    <section>
                        <h3>{self._escape_html(slide_content.title)}</h3>
                        <img src="{img_src}" style="max-height: 500px; max-width: 100%;" />
                    </section>
                    """
                else:
                    html = f"""
                    <section>
                        <h3>{self._escape_html(slide_content.title)}</h3>
                        <p>[Image not available]</p>
                    </section>
                    """
            else:
                html = ""

            slides_html.append(html)

        # Full HTML document
        html_content = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{self._escape_html(title)}</title>
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.5.0/reset.min.css">
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.5.0/reveal.min.css">
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.5.0/theme/{revealjs_theme}.min.css">
    <style>
        .reveal h1, .reveal h2, .reveal h3 {{
            text-transform: none;
        }}
        .reveal ul {{
            display: block;
        }}
        .reveal li {{
            margin: 0.5em 0;
        }}
    </style>
</head>
<body>
    <div class="reveal">
        <div class="slides">
            {"".join(slides_html)}
        </div>
    </div>
    <script src="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.5.0/reveal.min.js"></script>
    <script>
        Reveal.initialize({{
            hash: true,
            slideNumber: true,
            transition: 'slide'
        }});
    </script>
</body>
</html>
"""

        # Create ZIP package
        output_dir = Path(self.workspace_base) / "presentations"
        output_dir.mkdir(parents=True, exist_ok=True)

        zip_filename = f"presentation_{job_id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.zip"
        zip_path = output_dir / zip_filename

        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            zf.writestr('index.html', html_content)
            zf.writestr('README.txt', f"""
{title}
{'=' * len(title)}

This presentation was generated with LangConfig.

To view:
1. Extract this ZIP file
2. Open index.html in a web browser
3. Use arrow keys or click to navigate slides

The presentation uses Reveal.js (loaded from CDN) and requires an internet connection.
""")

        logger.info(f"Generated Reveal.js package at {zip_path}")
        return str(zip_path)

    async def _generate_google_slides(
        self,
        db: Session,
        title: str,
        slides: List[SlideContent],
        theme: str
    ) -> str:
        """
        Generate a Google Slides presentation.

        Args:
            db: Database session for OAuth credentials
            title: Presentation title
            slides: List of slide content
            theme: Visual theme

        Returns:
            Tuple of (URL to the created Google Slides presentation, list of image failure messages)
        """
        from googleapiclient.discovery import build
        from googleapiclient.http import MediaIoBaseUpload

        # Get valid credentials
        credentials = await google_oauth_service.get_valid_credentials(db)
        if not credentials:
            raise ValueError("Google OAuth not connected. Please connect your Google account first.")

        # Build Slides API client with HTTP timeout
        if httplib2 is not None and AuthorizedHttp is not None:
            http = httplib2.Http(timeout=60)
            authorized_http = AuthorizedHttp(credentials, http=http)
            slides_service = build('slides', 'v1', http=authorized_http)
            drive_service = build('drive', 'v3', http=authorized_http)
        else:
            slides_service = build('slides', 'v1', credentials=credentials)
            drive_service = build('drive', 'v3', credentials=credentials)

        try:
            # Create presentation
            presentation = await _execute(
                slides_service.presentations().create(body={'title': title})
            )
            presentation_id = presentation.get('presentationId')

            # Get the default slide (first slide created automatically)
            slides_response = await _execute(
                slides_service.presentations().get(presentationId=presentation_id)
            )
            default_slide_id = slides_response['slides'][0]['objectId']

            # Build batch update requests
            requests = []

            # Delete the default blank slide (we'll add our own)
            requests.append({
                'deleteObject': {
                    'objectId': default_slide_id
                }
            })

            # Add slides
            for i, slide_content in enumerate(slides):
                slide_id = f'slide_{i}'

                # Create slide
                if slide_content.slide_type == "title":
                    requests.append({
                        'createSlide': {
                            'objectId': slide_id,
                            'slideLayoutReference': {
                                'predefinedLayout': 'TITLE'
                            }
                        }
                    })
                elif slide_content.slide_type == "section":
                    requests.append({
                        'createSlide': {
                            'objectId': slide_id,
                            'slideLayoutReference': {
                                'predefinedLayout': 'SECTION_HEADER'
                            }
                        }
                    })
                else:
                    requests.append({
                        'createSlide': {
                            'objectId': slide_id,
                            'slideLayoutReference': {
                                'predefinedLayout': 'TITLE_AND_BODY'
                            }
                        }
                    })

            # Execute slide creation
            if requests:
                await _execute(
                    slides_service.presentations().batchUpdate(
                        presentationId=presentation_id,
                        body={'requests': requests}
                    )
                )

            # Get updated presentation to find placeholder IDs
            updated_presentation = await _execute(
                slides_service.presentations().get(presentationId=presentation_id)
            )

            # Build lookup by objectId for safe slide access
            slides_by_id = {s['objectId']: s for s in updated_presentation['slides']}

            # Add content to slides
            content_requests = []
            for i, slide_content in enumerate(slides):
                slide = slides_by_id.get(f'slide_{i}')
                if not slide:
                    logger.warning(f"Slide slide_{i} not found in presentation, skipping content")
                    continue
                slide_id = slide['objectId']

                # Find placeholders
                title_placeholder = None
                body_placeholder = None

                for element in slide.get('pageElements', []):
                    shape = element.get('shape', {})
                    placeholder = shape.get('placeholder', {})
                    placeholder_type = placeholder.get('type', '')

                    if placeholder_type in ['TITLE', 'CENTERED_TITLE']:
                        title_placeholder = element['objectId']
                    elif placeholder_type in ['BODY', 'SUBTITLE']:
                        body_placeholder = element['objectId']

                # Insert title text
                if title_placeholder and slide_content.title:
                    content_requests.append({
                        'insertText': {
                            'objectId': title_placeholder,
                            'text': slide_content.title
                        }
                    })

                # Insert body/bullets
                if body_placeholder:
                    if slide_content.subtitle:
                        content_requests.append({
                            'insertText': {
                                'objectId': body_placeholder,
                                'text': slide_content.subtitle
                            }
                        })
                    elif slide_content.bullets:
                        bullets_text = '\n'.join(f'• {b}' for b in slide_content.bullets)
                        content_requests.append({
                            'insertText': {
                                'objectId': body_placeholder,
                                'text': bullets_text
                            }
                        })

            # Execute content updates
            if content_requests:
                await _execute(
                    slides_service.presentations().batchUpdate(
                        presentationId=presentation_id,
                        body={'requests': content_requests}
                    )
                )

            # Handle images - upload to Drive and insert
            image_failures = []
            for i, slide_content in enumerate(slides):
                if slide_content.slide_type == "image" and slide_content.image_data:
                    try:
                        slide = slides_by_id.get(f'slide_{i}')
                        if not slide:
                            image_failures.append(f"Image {i}: slide not found")
                            continue
                        slide_id = slide['objectId']

                        # Upload image to Drive
                        image_bytes = base64.b64decode(slide_content.image_data)
                        media = MediaIoBaseUpload(
                            io.BytesIO(image_bytes),
                            mimetype=slide_content.image_mime_type,
                            resumable=True
                        )

                        file_metadata = {
                            'name': f'presentation_image_{i}.png',
                            'mimeType': slide_content.image_mime_type
                        }

                        uploaded_file = await _execute(
                            drive_service.files().create(
                                body=file_metadata,
                                media_body=media,
                                fields='id, webContentLink'
                            )
                        )

                        # Make file publicly accessible
                        await _execute(
                            drive_service.permissions().create(
                                fileId=uploaded_file['id'],
                                body={'type': 'anyone', 'role': 'reader'}
                            )
                        )

                        # Get direct download link
                        file_info = await _execute(
                            drive_service.files().get(
                                fileId=uploaded_file['id'],
                                fields='webContentLink'
                            )
                        )

                        image_url = file_info.get('webContentLink', '').replace('&export=download', '')

                        # Add image to slide
                        image_requests = [{
                            'createImage': {
                                'url': image_url,
                                'elementProperties': {
                                    'pageObjectId': slide_id,
                                    'size': {
                                        'width': {'magnitude': 500, 'unit': 'PT'},
                                        'height': {'magnitude': 350, 'unit': 'PT'}
                                    },
                                    'transform': {
                                        'scaleX': 1,
                                        'scaleY': 1,
                                        'translateX': 100,
                                        'translateY': 120,
                                        'unit': 'PT'
                                    }
                                }
                            }
                        }]

                        await _execute(
                            slides_service.presentations().batchUpdate(
                                presentationId=presentation_id,
                                body={'requests': image_requests}
                            )
                        )

                    except Exception as e:
                        logger.warning(f"Failed to add image to Google Slide: {e}")
                        image_failures.append(f"Image {i}: {e}")

            # Log partial image failures
            if image_failures:
                logger.warning(f"Some images failed to upload: {image_failures}")

            # Return the presentation URL
            presentation_url = f"https://docs.google.com/presentation/d/{presentation_id}/edit"
            logger.info(f"Created Google Slides presentation: {presentation_url}")

            return presentation_url, image_failures

        except Exception as e:
            logger.error(f"Failed to create Google Slides: {e}")
            raise

    def _escape_html(self, text: str) -> str:
        """Escape HTML special characters."""
        if not text:
            return ""
        return (
            text
            .replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
            .replace('"', "&quot;")
            .replace("'", "&#39;")
        )

    async def get_job(self, db: Session, job_id: int) -> Optional[PresentationJob]:
        """Get a presentation job by ID."""
        return db.query(PresentationJob).filter(PresentationJob.id == job_id).first()

    async def get_download_path(self, db: Session, job_id: int) -> Optional[str]:
        """Get the download path for a completed job."""
        job = await self.get_job(db, job_id)
        if not job or job.status != PresentationJobStatus.COMPLETED.value:
            return None
        return job.result_file_path


# Global instance
presentation_service = PresentationService()
